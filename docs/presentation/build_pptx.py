"""Build the BAI Core Trader / Hermes client presentation as .pptx.

Run:  python docs/presentation/build_pptx.py
Out:  docs/presentation/Hermes_BAI_Core.pptx

The deck is structured for a non-technical product owner: every slide
explains one concrete part of the bot - what it is, what it's
responsible for, and how it shows up in the UI. Marble/gold palette
mirrors the in-app theme.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).with_name("Hermes_BAI_Core.pptx")

# ── Palette ────────────────────────────────────────────────────────────
MARBLE = RGBColor(0xFB, 0xF7, 0xEC)
ALABASTER = RGBColor(0xF7, 0xEF, 0xD9)
PARCHMENT = RGBColor(0xEF, 0xE3, 0xC0)
GOLD = RGBColor(0xA8, 0x88, 0x4F)
GOLD_DEEP = RGBColor(0x7A, 0x5A, 0x1E)
GOLD_LIGHT = RGBColor(0xF4, 0xD7, 0x83)
NAVY = RGBColor(0x1B, 0x29, 0x40)
INK = RGBColor(0x3A, 0x33, 0x28)
MUTED = RGBColor(0x7E, 0x71, 0x63)
LAUREL = RGBColor(0x4C, 0x6B, 0x4A)
WINE = RGBColor(0x6B, 0x2A, 0x36)
BRONZE = RGBColor(0x9B, 0x7B, 0x2D)


prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


# ── Drawing helpers ────────────────────────────────────────────────────


def add_bg(slide, color=MARBLE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             font="Calibri", anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_rule(slide, x, y, w, color=GOLD, height=Pt(2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_card(slide, x, y, w, h, *, fill=ALABASTER, border=GOLD, border_w=0.75):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_w)
    card.shadow.inherit = False
    return card


def add_eyebrow(slide, x, y, text):
    add_text(slide, x, y, Inches(8), Inches(0.3),
             text.upper(), size=10, bold=True, color=GOLD_DEEP,
             font="Consolas")


def add_title(slide, x, y, w, text, *, size=36, color=NAVY):
    add_text(slide, x, y, w, Inches(1.2), text,
             size=size, bold=True, color=color, font="Cambria")


def add_subtitle(slide, x, y, w, text, *, size=15, color=MUTED):
    add_text(slide, x, y, w, Inches(0.8), text,
             size=size, color=color, font="Calibri", line_spacing=1.3)


def add_pill(slide, x, y, text, *, fill=GOLD_LIGHT, fg=GOLD_DEEP, width=None):
    if width is None:
        width = Inches(max(1.3, 0.22 + 0.085 * len(text)))
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, Inches(0.35))
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(9)
    r.font.bold = True
    r.font.name = "Consolas"
    r.font.color.rgb = fg
    return pill


def add_footer(slide, page, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             "BAI Core · Hermes Trading Bot · v1.0.40",
             size=9, color=MUTED, font="Consolas")
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{page} / {total}", size=9, color=MUTED, font="Consolas",
             align=PP_ALIGN.RIGHT)


# ── Slide builders ─────────────────────────────────────────────────────


SLIDES: list[callable] = []


def slide(fn):
    SLIDES.append(fn)
    return fn


@slide
def s_cover(prs, page, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, MARBLE)

    # Decorative side band
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              0, 0, Inches(0.35), SH)
    band.line.fill.background(); band.fill.solid()
    band.fill.fore_color.rgb = GOLD_DEEP
    band.shadow.inherit = False

    add_eyebrow(s, Inches(0.9), Inches(0.9), "BAI Core · Trading Intelligence")

    add_text(s, Inches(0.9), Inches(1.4), Inches(11), Inches(2),
             "Hermes", size=82, bold=True, color=NAVY, font="Cambria")
    add_text(s, Inches(0.9), Inches(3.0), Inches(11), Inches(1),
             "Бог торговли. Ваш бот.",
             size=24, color=GOLD_DEEP, font="Cambria")
    add_rule(s, Inches(0.9), Inches(3.85), Inches(2.5), color=GOLD)

    add_text(s, Inches(0.9), Inches(4.05), Inches(11), Inches(2),
             "Полностью автономный десктоп-бот для MetaTrader 5.\n"
             "Аналитика, риск-менеджмент и исполнение - на вашем "
             "компьютере, без облака.",
             size=16, color=MUTED, font="Calibri", line_spacing=1.4)

    add_text(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.4),
             "Презентация функционала · версия 1.0.40 · май 2026",
             size=10, color=MUTED, font="Consolas")


@slide
def s_problem(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Шаг 0 · Зачем нужен Hermes")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Ручная торговля - это три проблемы сразу")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    problems = [
        ("Эмоции", "Страх и жадность ломают любую стратегию. "
                   "Трейдер видит просадку - закрывает рано. Видит "
                   "профит - держит слишком долго.", WINE),
        ("Время", "Рынок работает 24/5. Человек физически не может "
                  "анализировать 6 пар каждую минуту. Сигналы "
                  "пропускаются, точки входа уходят.", BRONZE),
        ("Дисциплина", "SL/TP, дневной лимит сделок, размер позиции - "
                       "правила есть у всех. Соблюдают единицы. Один "
                       "срыв - и месяц работы насмарку.", LAUREL),
    ]
    col_w = Inches(4.0)
    col_gap = Inches(0.25)
    start_x = Inches(0.6)
    for i, (title, body, tone) in enumerate(problems):
        x = start_x + (col_w + col_gap) * i
        add_card(s, x, Inches(2.5), col_w, Inches(3.6), fill=ALABASTER, border=tone)
        add_pill(s, x + Inches(0.3), Inches(2.7), title.upper(),
                 fill=tone, fg=MARBLE)
        add_text(s, x + Inches(0.3), Inches(3.3), col_w - Inches(0.6),
                 Inches(2.5), body, size=14, color=INK, line_spacing=1.35)

    add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.4),
             "Hermes решает все три - спокойно, систематически, без выходных.",
             size=14, bold=True, color=NAVY, font="Cambria")
    add_footer(s, page, total)


@slide
def s_what(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Что такое Hermes")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Десктоп-приложение, которое торгует за вас")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    items = [
        ("📦 Один установщик", "Hermes-Setup.exe двойным кликом. Никаких "
                                "облачных аккаунтов, никакой регистрации. "
                                "Бот работает офлайн на вашем ПК."),
        ("🔐 MT5 + криптобиржи", "MetaTrader 5 уже подключён. CCXT-адаптер "
                                  "для Binance/Bybit готов, активируется в "
                                  "одну строку конфига."),
        ("🧠 Объяснимые решения", "Бот пишет под каждой сделкой ПОЧЕМУ он её "
                                   "открыл - какие индикаторы согласились, "
                                   "какая стратегия сработала, какая уверенность."),
        ("🛡️ 3 уровня защиты", "Risk Engine следит за просадкой, "
                                "брокерский SL/TP уходит на сервер при открытии, "
                                "Kill-switch закрывает всё одной кнопкой."),
        ("📊 Живой дашборд", "Эквити-кривая, открытые позиции, P&L по "
                              "режимам, состояние пар - всё обновляется в "
                              "реальном времени через WebSocket."),
        ("🔁 Самокалибровка", "Раз в неделю бот пересматривает параметры "
                               "стратегий на свежей истории. Walk-forward, "
                               "автоматический откат если стало хуже."),
    ]
    col_w = Inches(6.0)
    col_h = Inches(1.5)
    for i, (title, body) in enumerate(items):
        col, row = i % 2, i // 2
        x = Inches(0.6) + col * (col_w + Inches(0.25))
        y = Inches(2.45) + row * (col_h + Inches(0.15))
        add_card(s, x, y, col_w, col_h, fill=ALABASTER, border=GOLD)
        add_text(s, x + Inches(0.25), y + Inches(0.15),
                 col_w - Inches(0.5), Inches(0.45),
                 title, size=14, bold=True, color=NAVY, font="Cambria")
        add_text(s, x + Inches(0.25), y + Inches(0.6),
                 col_w - Inches(0.5), Inches(0.9),
                 body, size=11, color=INK, line_spacing=1.3)
    add_footer(s, page, total)


@slide
def s_arch(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Архитектура")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Один процесс, четыре слоя ответственности")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    layers = [
        ("Frontend (UI)", "React 18 + Vite + TanStack Query.\n"
                          "Дашборд, сделки, стратегия, бэктест, "
                          "адаптация, уведомления, настройки.",
         "Что видит и нажимает оператор."),
        ("Backend (API)", "FastAPI + async SQLAlchemy.\n"
                          "REST + WebSocket. Auth, trading, brokers, "
                          "positions, trades, adaptive, system.",
         "Мозг приложения. Принимает команды, отдаёт данные."),
        ("Trading Engine", "Worker + StrategyRunner + IndicatorPanel + "
                           "RiskEngine + StrategyEnsemble.\n"
                           "Раз в 60с: анализ, риск-чек, ордер.",
         "Алгоритм торговли как таковой."),
        ("Brokers", "MT5 (нативный SDK) и CCXT (крипта). "
                    "Watchdog с авто-переподключением. "
                    "Vault для хранения паролей.",
         "Интерфейс к внешнему миру: котировки, ордера."),
    ]
    row_h = Inches(1.05)
    for i, (name, what, owns) in enumerate(layers):
        y = Inches(2.4) + i * (row_h + Inches(0.1))
        add_card(s, Inches(0.6), y, Inches(12.1), row_h)
        add_pill(s, Inches(0.85), y + Inches(0.32),
                 f"СЛОЙ {i+1}", fill=GOLD_DEEP, fg=MARBLE, width=Inches(0.95))
        add_text(s, Inches(1.95), y + Inches(0.1), Inches(3.5), Inches(0.45),
                 name, size=15, bold=True, color=NAVY, font="Cambria")
        add_text(s, Inches(1.95), y + Inches(0.55), Inches(3.5),
                 Inches(0.55), what, size=10, color=INK, line_spacing=1.25)
        add_text(s, Inches(5.7), y + Inches(0.32), Inches(7.0), Inches(0.45),
                 f"Ответственность: {owns}",
                 size=12, color=GOLD_DEEP, font="Cambria", line_spacing=1.3)
    add_footer(s, page, total)


@slide
def s_modes(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Режимы работы")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Три способа торговли - оператор выбирает один")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    modes = [
        ("ПРОВЕРЕННЫЙ", LAUREL,
         "Одна стратегия, 3-5 пар",
         [("Стратегии", "Только TrendFollowing"),
          ("Пары", "Из конфига (3-5 мажоров)"),
          ("Порог уверенности", "≥ 0.70 (строго)"),
          ("Лимит сделок", "1-3 в день"),
          ("Лот", "0.5% от эквити (мин 0.01)"),
          ("Для кого", "Консервативно, ровная просадка")]),
        ("АВТОНОМНЫЙ", GOLD_DEEP,
         "Все стратегии, любая пара",
         [("Стратегии", "Trend + Momentum (по умолчанию)"),
          ("Пары", "Все включённые в конфиге"),
          ("Порог уверенности", "≥ 0.50 (умеренно)"),
          ("Лимит сделок", "1-3 в день"),
          ("Лот", "0.5% от эквити (мин 0.01)"),
          ("Для кого", "Агрессивнее, ищет лучший вход")]),
        ("ТЕСТ-СДЕЛКА", BRONZE,
         "Разовое открытие для проверки",
         [("Стратегии", "Полный ансамбль (как Автономный)"),
          ("Пары", "Все из конфига - бот выбирает лучшую"),
          ("Порог уверенности", "≥ 0.40"),
          ("Лимит сделок", "Один разовый клик"),
          ("Лот", "0.5% от эквити"),
          ("Для кого", "Проверить связь с брокером")]),
    ]
    col_w = Inches(4.0)
    for i, (name, tone, sub, rows) in enumerate(modes):
        x = Inches(0.6) + i * (col_w + Inches(0.25))
        y = Inches(2.5)
        h = Inches(4.2)
        add_card(s, x, y, col_w, h, fill=ALABASTER, border=tone, border_w=1.2)
        add_text(s, x + Inches(0.3), y + Inches(0.2),
                 col_w - Inches(0.6), Inches(0.4),
                 name, size=15, bold=True, color=tone, font="Cambria")
        add_text(s, x + Inches(0.3), y + Inches(0.6),
                 col_w - Inches(0.6), Inches(0.3),
                 sub, size=10, color=MUTED, font="Calibri")
        add_rule(s, x + Inches(0.3), y + Inches(0.95),
                 col_w - Inches(0.6), color=tone, height=Pt(0.5))
        for j, (k, v) in enumerate(rows):
            ry = y + Inches(1.05) + j * Inches(0.5)
            add_text(s, x + Inches(0.3), ry,
                     col_w - Inches(0.6), Inches(0.22),
                     k.upper(),
                     size=8, bold=True, color=GOLD_DEEP, font="Consolas")
            add_text(s, x + Inches(0.3), ry + Inches(0.22),
                     col_w - Inches(0.6), Inches(0.28),
                     v, size=10, color=INK)
    add_footer(s, page, total)


@slide
def s_tick(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Что делает бот каждые 60 секунд")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Цикл одного tick'а - семь шагов подряд")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    steps = [
        ("1", "Health probe",
         "Проверка что MT5 живой. 3 фейла подряд - пауза + уведомление."),
        ("2", "OHLCV для каждой пары",
         "Скачиваем 300 свечей по таймфрейму конфига для всех символов."),
        ("3", "IndicatorPanel",
         "Считаем 19 индикаторов: RSI, MACD, BB, ATR, EMA 50/200, "
         "ADX, +DI/-DI, Stochastic, Donchian, ROC и другие."),
        ("4", "Strategy Ensemble",
         "4 стратегии формируют сигналы Long/Short/Flat с обоснованием. "
         "Голосование (majority/any/all) → одно решение по символу."),
        ("5", "Risk Engine",
         "Проверка: дневной убыток < 5%, drawdown < 10%, "
         "позиций ≤ 5. Если нарушено - вход запрещён."),
        ("6", "Order placement",
         "Выбирается лучший сигнал по уверенности. Лот = 0.5% эквити, "
         "снап под volume_step. SL/TP уходят на сервер брокера."),
        ("7", "Snapshot + Broadcast",
         "EquityPoint + PositionSnapshot пишутся в SQLite. "
         "WebSocket рассылает обновления всем подключённым клиентам."),
    ]
    for i, (n, title, body) in enumerate(steps):
        col = i // 4 if i < 7 else 1
        row = i % 4 if i < 4 else (i - 4)
        x = Inches(0.6) + col * Inches(6.2)
        y = Inches(2.4) + row * Inches(1.1)
        # circle index
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.55), Inches(0.55))
        circ.fill.solid(); circ.fill.fore_color.rgb = GOLD_DEEP
        circ.line.fill.background()
        tf = circ.text_frame; tf.margin_left = tf.margin_right = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = n
        r.font.size = Pt(18); r.font.bold = True
        r.font.color.rgb = MARBLE; r.font.name = "Cambria"
        # title + body
        add_text(s, x + Inches(0.75), y, Inches(5.2), Inches(0.4),
                 title, size=14, bold=True, color=NAVY, font="Cambria")
        add_text(s, x + Inches(0.75), y + Inches(0.42), Inches(5.2),
                 Inches(0.7), body, size=10, color=INK, line_spacing=1.3)
    add_footer(s, page, total)


@slide
def s_indicators(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Indicator Panel")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "19 индикаторов на каждой паре каждый тик")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    groups = [
        ("МОМЕНТУМ", [
            "RSI (Wilder)", "Stochastic %K", "Stochastic %D",
            "ROC (Rate of Change)", "MACD line", "MACD signal", "MACD hist",
        ]),
        ("ТРЕНД", [
            "EMA 50 (быстрая)", "EMA 200 (медленная)",
            "ADX (сила тренда)", "+DI (покупатели)", "-DI (продавцы)",
        ]),
        ("ВОЛАТИЛЬНОСТЬ", [
            "ATR (абсолют)", "ATR% (нормированный)",
            "BB верхняя", "BB средняя", "BB нижняя",
        ]),
        ("ПРОБОИ", [
            "Donchian верхняя", "Donchian нижняя",
        ]),
    ]
    col_w = Inches(2.95)
    for i, (group, items) in enumerate(groups):
        x = Inches(0.6) + i * (col_w + Inches(0.15))
        add_card(s, x, Inches(2.4), col_w, Inches(4.4))
        add_pill(s, x + Inches(0.25), Inches(2.6), group,
                 fill=GOLD_DEEP, fg=MARBLE,
                 width=Inches(min(2.5, 0.3 + 0.085 * len(group))))
        for j, it in enumerate(items):
            ity = Inches(3.15) + j * Inches(0.42)
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.3), ity + Inches(0.13),
                                     Inches(0.08), Inches(0.08))
            dot.fill.solid(); dot.fill.fore_color.rgb = GOLD
            dot.line.fill.background()
            add_text(s, x + Inches(0.5), ity, col_w - Inches(0.7),
                     Inches(0.35), it, size=11, color=INK)

    add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
             "Стандартные периоды по умолчанию, всё конфигурируется через UI · "
             "значения транслируются в /ws/signals для дашборда",
             size=10, color=MUTED, font="Calibri")
    add_footer(s, page, total)


@slide
def s_strategies(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Strategy Ensemble")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "4 независимые стратегии голосуют большинством")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    strategies = [
        ("Trend Following", LAUREL,
         "EMA50 vs EMA200 + ADX > 25",
         "Входит когда EMA быстрая выше/ниже медленной "
         "И тренд устойчивый (ADX). Подтверждение через +DI/-DI."),
        ("Mean Reversion", BRONZE,
         "BB-касание + RSI экстремум",
         "Цена коснулась границы Боллинджера И RSI в зоне "
         "перепроданности/перекупленности. Подтверждение Stochastic-разворотом."),
        ("Breakout", WINE,
         "Пробой Donchian + ATR-фильтр",
         "Цена пробила канал Дончиана с запасом 0.05%. "
         "Только если ATR > 0.08% (фильтр от мёртвого рынка)."),
        ("Momentum", GOLD_DEEP,
         "MACD-гистограмма + Stochastic",
         "MACD-гистограмма положительна/отрицательна И Stochastic "
         "пересекает в нужном направлении ниже/выше 50."),
    ]
    col_w = Inches(6.0)
    for i, (name, tone, formula, body) in enumerate(strategies):
        col, row = i % 2, i // 2
        x = Inches(0.6) + col * (col_w + Inches(0.25))
        y = Inches(2.45) + row * Inches(1.95)
        add_card(s, x, y, col_w, Inches(1.8))
        add_pill(s, x + Inches(0.25), y + Inches(0.2),
                 name.upper(), fill=tone, fg=MARBLE,
                 width=Inches(min(3.0, 0.3 + 0.11 * len(name))))
        add_text(s, x + Inches(0.25), y + Inches(0.7),
                 col_w - Inches(0.5), Inches(0.3),
                 formula, size=11, bold=True, color=tone, font="Consolas")
        add_text(s, x + Inches(0.25), y + Inches(1.0),
                 col_w - Inches(0.5), Inches(0.7),
                 body, size=11, color=INK, line_spacing=1.3)

    add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             "Голосование majority: ≥50% активных стратегий должны "
             "согласиться. Иначе - бот пропускает тик.",
             size=11, color=MUTED, font="Calibri")
    add_footer(s, page, total)


@slide
def s_risk(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Risk Engine")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Три независимых circuit breaker'а")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    breakers = [
        ("ДНЕВНОЙ УБЫТОК", "-5%", LAUREL,
         "Если эквити упало больше 5% от стартового за сессию - "
         "бот перестаёт открывать новые позиции. Старые остаются "
         "под защитой брокерских SL.",
         "Сбрасывается в полночь UTC."),
        ("DRAWDOWN", "-10%", BRONZE,
         "От пикового эквити сессии. Защита от затяжной просадки: "
         "если максимум был 50k и упали до 45k - стоп. Не даёт "
         "ансамблю усугубить положение.",
         "Сбрасывается только перезапуском."),
        ("MAX ПОЗИЦИЙ", "5", WINE,
         "Сколько открытых позиций может быть одновременно. "
         "Защита от перегрева портфеля - если 5 уже работают, "
         "новые входы блокируются до закрытия одной из них.",
         "Считается в реальном времени."),
    ]
    col_w = Inches(4.0)
    for i, (name, threshold, tone, body, footnote) in enumerate(breakers):
        x = Inches(0.6) + i * (col_w + Inches(0.25))
        y = Inches(2.4)
        h = Inches(4.4)
        add_card(s, x, y, col_w, h, border=tone, border_w=1.2)
        # Big threshold number
        add_text(s, x + Inches(0.3), y + Inches(0.2),
                 col_w - Inches(0.6), Inches(1.3),
                 threshold, size=54, bold=True, color=tone, font="Cambria",
                 align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.3), y + Inches(1.65),
                 col_w - Inches(0.6), Inches(0.4),
                 name, size=12, bold=True, color=GOLD_DEEP, font="Consolas",
                 align=PP_ALIGN.CENTER)
        add_rule(s, x + Inches(0.5), y + Inches(2.05),
                 col_w - Inches(1.0), color=tone, height=Pt(0.5))
        add_text(s, x + Inches(0.3), y + Inches(2.2),
                 col_w - Inches(0.6), Inches(1.6),
                 body, size=11, color=INK, line_spacing=1.35)
        add_text(s, x + Inches(0.3), y + Inches(3.8),
                 col_w - Inches(0.6), Inches(0.5),
                 footnote, size=10, color=MUTED, font="Consolas")

    add_footer(s, page, total)


@slide
def s_orders(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Управление ордерами")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "От сигнала до исполнения - что происходит с ордером")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    flow = [
        ("Сигнал ансамбля",
         "Стратегии согласились: LONG USDJPY, уверенность 0.62."),
        ("Расчёт лота",
         "0.5% × эквити / 100k = ~0.0025. Снапаем под "
         "volume_step символа = 0.01."),
        ("Round to step + clamp",
         "round(volume/step)*step. Зажимаем между volume_min и "
         "volume_max - брокер не отвергнет с 'Invalid volume'."),
        ("Filling mode",
         "Читаем symbol_info.filling_mode (битмаска). "
         "FOK → IOC → RETURN по предпочтению."),
        ("SL/TP на брокера",
         "2×ATR стоп, 4×ATR таргет (1:2 R/R). "
         "Уходят как абсолютные цены вместе с ордером."),
        ("order_send",
         "MT5 принимает заявку. Retcode TRADE_RETCODE_DONE = успех. "
         "Любой другой - переводим в человеческое сообщение."),
        ("Persist + broadcast",
         "TradeRow в SQLite с обоснованием стратегии. "
         "WebSocket /ws/signals: 'trade_opened' для UI и Telegram."),
    ]
    item_h = Inches(0.55)
    for i, (title, body) in enumerate(flow):
        y = Inches(2.45) + i * (item_h + Inches(0.07))
        # Number disk
        n = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(0.7), y + Inches(0.1),
                               Inches(0.4), Inches(0.4))
        n.fill.solid(); n.fill.fore_color.rgb = GOLD_DEEP
        n.line.fill.background()
        tf = n.text_frame; tf.margin_left = tf.margin_right = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(i + 1)
        r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = MARBLE; r.font.name = "Cambria"
        # Card
        add_card(s, Inches(1.3), y, Inches(11.4), item_h, fill=ALABASTER)
        add_text(s, Inches(1.55), y + Inches(0.1), Inches(3.3), Inches(0.4),
                 title, size=13, bold=True, color=NAVY, font="Cambria")
        add_text(s, Inches(4.95), y + Inches(0.13), Inches(7.6),
                 Inches(0.4), body, size=11, color=INK, line_spacing=1.2)
    add_footer(s, page, total)


@slide
def s_explainable(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Объяснимость")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Каждая сделка приходит с обоснованием")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    # Mock UI card
    add_card(s, Inches(0.6), Inches(2.4), Inches(12.1), Inches(4.4))
    add_pill(s, Inches(0.85), Inches(2.6), "USDJPY · LONG · 0.01 · #285741",
             fill=LAUREL, fg=MARBLE, width=Inches(4.4))
    add_text(s, Inches(0.85), Inches(3.1), Inches(11.5), Inches(0.5),
             "LONG по USDJPY: 2 из 2 стратегий согласны "
             "(средняя уверенность 0.62).",
             size=15, bold=True, color=NAVY, font="Cambria")
    bullets = [
        ("Trend Following",
         "Быстрая EMA (148.45) выше медленной (147.80) - тренд "
         "восходящий. ADX 33.5 (порог 25) подтверждает силу. "
         "+DI 28.1 превышает -DI 17.4 - покупатели доминируют."),
        ("Momentum",
         "Гистограмма MACD +0.00527 положительна - импульс растёт. "
         "Stochastic %K 64.1 пересёк %D 61.4 снизу вверх. "
         "Движение наверх набирает обороты."),
    ]
    y0 = Inches(3.95)
    for i, (name, body) in enumerate(bullets):
        y = y0 + i * Inches(1.25)
        bul = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(1.0), y + Inches(0.18),
                                 Inches(0.13), Inches(0.13))
        bul.fill.solid(); bul.fill.fore_color.rgb = GOLD
        bul.line.fill.background()
        add_text(s, Inches(1.3), y, Inches(11.0), Inches(0.4),
                 name, size=13, bold=True, color=GOLD_DEEP, font="Cambria")
        add_text(s, Inches(1.3), y + Inches(0.4), Inches(11.0), Inches(0.9),
                 body, size=11, color=INK, line_spacing=1.35)

    add_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
             "Видно в Журнале сделок · хранится в БД · можно добавить "
             "свою заметку для разбора post-mortem",
             size=10, color=MUTED, font="Calibri")
    add_footer(s, page, total)


@slide
def s_ui(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Интерфейс оператора")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "11 экранов - всё что нужно для контроля")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    screens = [
        ("🏛️ Дашборд", "Баланс, equity-кривая, P&L по режимам, "
                       "открытые позиции, drawdown, состояние пар"),
        ("📋 Сделки", "История с обоснованием каждой. "
                     "Заметки оператора, фильтры, CSV-экспорт"),
        ("⚙️ Стратегия", "Выбор активных стратегий, "
                        "пары, таймфрейм, ансамбль, режим голосования"),
        ("🔌 Брокеры", "MT5/CCXT подключения, "
                      "тест связи, статус, переподключение"),
        ("🧪 Бэктест", "Прогон на исторических данных, "
                      "Sharpe/winrate/drawdown, equity-кривая"),
        ("🎚️ Оптимизация", "Grid search по параметрам, "
                          "лидерборд, применение лучшего"),
        ("🌡️ Адаптация", "Регимы пар (trend/flat/high_vol), "
                         "walk-forward, история калибровок"),
        ("📱 На телефон", "QR-код для PWA, "
                         "статус ngrok-туннеля если включён"),
        ("🔔 Уведомления", "Telegram-бот, web-push, "
                          "выбор событий, тест-сообщение"),
        ("📜 Журнал", "Лог backend'а, "
                     "скачать архив для поддержки"),
        ("🛠️ Настройки", "Тема, обновления, переход на v+1, "
                        "сброс тура, ссылка на лог"),
    ]
    col_w = Inches(4.0)
    for i, (icon_title, body) in enumerate(screens):
        col = i % 3
        row = i // 3
        x = Inches(0.6) + col * (col_w + Inches(0.15))
        y = Inches(2.4) + row * Inches(1.05)
        add_card(s, x, y, col_w, Inches(0.95), fill=ALABASTER)
        add_text(s, x + Inches(0.2), y + Inches(0.1),
                 col_w - Inches(0.4), Inches(0.35),
                 icon_title, size=12, bold=True, color=NAVY, font="Cambria")
        add_text(s, x + Inches(0.2), y + Inches(0.45),
                 col_w - Inches(0.4), Inches(0.45),
                 body, size=9, color=INK, line_spacing=1.25)
    add_footer(s, page, total)


@slide
def s_adaptive(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Адаптивная калибровка")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Бот сам подстраивает параметры под текущий рынок")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    blocks = [
        ("Регим-детектор",
         "Каждые 5 минут бот определяет режим каждой пары: "
         "trend / flat / high_vol. По ADX, ATR%, EMA-выравниванию, "
         "Hurst-экспоненте. Это видно в Дашборде как 'Состояние пар'."),
        ("Walk-forward",
         "Раз в неделю (по умолчанию) запускается переподбор: "
         "обучение на N-дневном окне, тест на (N+7)-дневном. "
         "Шкала Sharpe-ratio определяет победителя."),
        ("Challenger vs Champion",
         "Текущие параметры (champion) против переподобранных "
         "(challenger). Применяется только если Walk-forward Score "
         "новых выше. Без подтверждения - бот работает на старых."),
        ("Откат одной кнопкой",
         "Если применили новый параметр и за 24-48 часов P&L упал - "
         "в UI кнопка 'Откатить'. Возвращаются предыдущие значения. "
         "История всех калибровок хранится в БД."),
    ]
    col_w = Inches(6.0)
    for i, (title, body) in enumerate(blocks):
        col, row = i % 2, i // 2
        x = Inches(0.6) + col * (col_w + Inches(0.25))
        y = Inches(2.45) + row * Inches(2.05)
        add_card(s, x, y, col_w, Inches(1.85))
        add_text(s, x + Inches(0.3), y + Inches(0.2),
                 col_w - Inches(0.6), Inches(0.4),
                 title, size=14, bold=True, color=NAVY, font="Cambria")
        add_rule(s, x + Inches(0.3), y + Inches(0.65),
                 col_w - Inches(0.6), color=GOLD, height=Pt(0.5))
        add_text(s, x + Inches(0.3), y + Inches(0.75),
                 col_w - Inches(0.6), Inches(1.1),
                 body, size=11, color=INK, line_spacing=1.35)
    add_footer(s, page, total)


@slide
def s_notifications(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Уведомления")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Бот сообщает что происходит - даже когда вы не за компьютером")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    channels = [
        ("Telegram-бот", LAUREL,
         "Подключается одной строкой (BotFather token + chat_id). "
         "В UI - тест-сообщение. События: открытие сделки, kill-switch, "
         "broker_down, риск-блокировки."),
        ("Web-push", BRONZE,
         "Стандартные браузерные уведомления через Service Worker. "
         "Работает когда дашборд открыт в браузере или мобильной PWA. "
         "VAPID-ключи в настройках."),
        ("Native (Windows)", GOLD_DEEP,
         "Toast-уведомления Windows 10/11 через PyWebView мост. "
         "Активируются автоматически в десктоп-сборке. "
         "Полезно при работе с приложением в фоне."),
    ]
    col_w = Inches(4.0)
    for i, (name, tone, body) in enumerate(channels):
        x = Inches(0.6) + i * (col_w + Inches(0.25))
        y = Inches(2.4)
        add_card(s, x, y, col_w, Inches(3.5), border=tone, border_w=1.2)
        add_pill(s, x + Inches(0.3), y + Inches(0.25), name.upper(),
                 fill=tone, fg=MARBLE,
                 width=Inches(min(3.3, 0.3 + 0.11 * len(name))))
        add_text(s, x + Inches(0.3), y + Inches(0.9),
                 col_w - Inches(0.6), Inches(2.5),
                 body, size=12, color=INK, line_spacing=1.4)

    # Events list
    add_card(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.85))
    add_text(s, Inches(0.85), Inches(6.15), Inches(11.6), Inches(0.3),
             "СОБЫТИЯ КОТОРЫЕ ШЛЮТСЯ:",
             size=10, bold=True, color=GOLD_DEEP, font="Consolas")
    add_text(s, Inches(0.85), Inches(6.45), Inches(11.6), Inches(0.4),
             "Открытие сделки · закрытие по SL/TP · kill-switch · "
             "брокер недоступен · risk_block · ошибка стратегии · "
             "успешная калибровка",
             size=11, color=INK)
    add_footer(s, page, total)


@slide
def s_security(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Безопасность")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Пароли брокера никогда не покидают ваш ПК")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    items = [
        ("Vault (хранилище)",
         "Все секреты (пароль MT5, API-ключи, Telegram-токены) "
         "хранятся в зашифрованном файле на вашем диске."),
        ("Argon2id KDF",
         "Ключ шифрования выводится через Argon2id (победитель PHC) "
         "с памятью 64 МБ и 3 итерациями - устойчиво к GPU-атакам."),
        ("Fernet AES-128-CBC + HMAC",
         "Сам vault шифруется Fernet'ом: AES-128 в режиме CBC + "
         "HMAC-SHA256 для проверки целостности."),
        ("Атомарная запись",
         "Файл vault'а перезаписывается через tmp + rename - "
         "при крахе посередине сохранится либо старая, либо новая "
         "версия, никогда не битая."),
        ("Passwordless mode (v1.0.31+)",
         "Ключ выводится из неподдельных атрибутов установки. "
         "Vault разблокируется автоматически - оператору не нужно "
         "помнить мастер-пароль."),
        ("Локальный pip-only бэкенд",
         "FastAPI слушает только 127.0.0.1, эфемерный порт. "
         "Извне ПК недоступен - даже на той же Wi-Fi."),
    ]
    col_w = Inches(6.0)
    for i, (title, body) in enumerate(items):
        col, row = i % 2, i // 2
        x = Inches(0.6) + col * (col_w + Inches(0.25))
        y = Inches(2.4) + row * Inches(1.45)
        add_card(s, x, y, col_w, Inches(1.35))
        add_text(s, x + Inches(0.25), y + Inches(0.15),
                 col_w - Inches(0.5), Inches(0.4),
                 title, size=13, bold=True, color=NAVY, font="Cambria")
        add_text(s, x + Inches(0.25), y + Inches(0.55),
                 col_w - Inches(0.5), Inches(0.8),
                 body, size=10, color=INK, line_spacing=1.3)
    add_footer(s, page, total)


@slide
def s_stack(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Технологии")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Production-ready стек - всё что используется")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    blocks = [
        ("Backend", [
            "Python 3.13 · FastAPI · uvicorn (ASGI)",
            "SQLAlchemy 2 (async) · aiosqlite · Alembic",
            "MetaTrader5 SDK · CCXT · pandas · numpy",
            "APScheduler · websockets · httpx",
            "cryptography (Fernet, Argon2id)",
        ]),
        ("Frontend", [
            "React 18 · TypeScript · Vite",
            "TanStack Query · React Router 6",
            "TailwindCSS · lucide-react · Cambria/Calibri",
            "lightweight-charts (TradingView)",
            "vite-plugin-pwa (Service Worker)",
        ]),
        ("Desktop", [
            "PyWebView 5 · Edge WebView2 runtime",
            "PyInstaller onedir bundle",
            "Inno Setup installer (Windows)",
            "Auto-update механизм встроен",
            "Single-instance + log rotation",
        ]),
        ("DevOps", [
            "GitHub Actions release.yml",
            "Tag-driven version injection",
            "Hermes-Setup-{ver}.exe + portable.zip",
            "pytest (unit + e2e) · TypeScript strict",
            "SHA256SUMS для каждой сборки",
        ]),
    ]
    col_w = Inches(6.0)
    for i, (title, items) in enumerate(blocks):
        col, row = i % 2, i // 2
        x = Inches(0.6) + col * (col_w + Inches(0.25))
        y = Inches(2.4) + row * Inches(2.2)
        add_card(s, x, y, col_w, Inches(2.05))
        add_pill(s, x + Inches(0.25), y + Inches(0.2), title,
                 fill=GOLD_DEEP, fg=MARBLE,
                 width=Inches(min(2.5, 0.3 + 0.11 * len(title))))
        for j, it in enumerate(items):
            ity = y + Inches(0.7) + j * Inches(0.28)
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.3), ity + Inches(0.1),
                                     Inches(0.06), Inches(0.06))
            dot.fill.solid(); dot.fill.fore_color.rgb = GOLD
            dot.line.fill.background()
            add_text(s, x + Inches(0.5), ity, col_w - Inches(0.7),
                     Inches(0.28), it, size=10, color=INK)
    add_footer(s, page, total)


@slide
def s_install(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Установка")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Двойной клик - три минуты - готов к работе")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    steps = [
        ("1", "Скачать Hermes-Setup.exe",
         "С GitHub Releases - https://github.com/artsignstudiokz/Hermes. "
         "Размер ~110 МБ, подписан SHA256."),
        ("2", "Двойной клик - 'Установить'",
         "Inno Setup установщик. Прогресс-бар, лицензия, "
         "выбор папки. Без админских прав - в user-space."),
        ("3", "MetaTrader 5 должен быть запущен",
         "Hermes использует уже залогиненный терминал. "
         "Без MT5 - бот покажет 'Брокер недоступен'."),
        ("4", "Открыть Hermes - подключить брокера",
         "Брокеры → Добавить → MT5 → выбрать сервер из списка "
         "(Exness/IC Markets/etc.) → логин/пароль."),
        ("5", "Стратегия → Запустить",
         "Главная → Проверенный/Автономный. Бот стартует, "
         "анализирует пары, ищет точки входа."),
        ("6", "Обновления",
         "Hermes сам проверяет новые версии. Кнопка 'Обновить' "
         "в Настройках. Установщик новой версии ставится поверх, "
         "БД и vault сохраняются."),
    ]
    for i, (n, title, body) in enumerate(steps):
        col = i // 3
        row = i % 3
        x = Inches(0.6) + col * Inches(6.2)
        y = Inches(2.5) + row * Inches(1.45)
        # number disk
        ds = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.7), Inches(0.7))
        ds.fill.solid(); ds.fill.fore_color.rgb = GOLD_DEEP
        ds.line.fill.background()
        tf = ds.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Pt(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = n
        r.font.size = Pt(20); r.font.bold = True
        r.font.color.rgb = MARBLE; r.font.name = "Cambria"
        # texts
        add_text(s, x + Inches(0.9), y, Inches(5.1), Inches(0.45),
                 title, size=14, bold=True, color=NAVY, font="Cambria")
        add_text(s, x + Inches(0.9), y + Inches(0.45), Inches(5.1),
                 Inches(0.95), body, size=11, color=INK, line_spacing=1.3)
    add_footer(s, page, total)


@slide
def s_quality(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Качество кода")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Многоуровневая валидация на всём пути")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    metrics = [
        ("76", "Pytest-тестов", "unit + integration"),
        ("100%", "TypeScript strict", "noUncheckedIndexedAccess"),
        ("19", "Индикаторов", "технического анализа"),
        ("4", "Стратегии", "независимых ансамбля"),
        ("3", "Risk-контура", "независимых breaker'а"),
        ("60с", "Tick-интервал", "каждый цикл"),
        ("11", "Экранов", "в UI"),
        ("40", "Релизов", "за май 2026"),
    ]
    col_w = Inches(3.0)
    for i, (val, title, sub) in enumerate(metrics):
        col, row = i % 4, i // 4
        x = Inches(0.6) + col * (col_w + Inches(0.15))
        y = Inches(2.4) + row * Inches(2.0)
        add_card(s, x, y, col_w, Inches(1.85))
        add_text(s, x, y + Inches(0.2), col_w, Inches(0.9),
                 val, size=48, bold=True, color=GOLD_DEEP, font="Cambria",
                 align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.1), col_w, Inches(0.35),
                 title, size=12, bold=True, color=NAVY, font="Cambria",
                 align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(1.45), col_w, Inches(0.3),
                 sub, size=10, color=MUTED, font="Consolas",
                 align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.6), Inches(6.65), Inches(12), Inches(0.4),
             "CI прогоняет тесты на каждый коммит. "
             "Релизы пушатся только из main с зелёным CI.",
             size=11, color=MUTED, font="Calibri", align=PP_ALIGN.CENTER)
    add_footer(s, page, total)


@slide
def s_roadmap(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)
    add_eyebrow(s, Inches(0.6), Inches(0.5), "Дорожная карта")
    add_title(s, Inches(0.6), Inches(0.85), Inches(12),
              "Что уже работает · что ближайшее")
    add_rule(s, Inches(0.6), Inches(2.0), Inches(2.5))

    done = [
        "Дашборд с live-данными (WebSocket)",
        "MT5 + watchdog с авто-переподключением",
        "4 стратегии × 19 индикаторов · объяснимые сигналы",
        "Risk Engine (3 контура) + брокерские SL/TP",
        "Два режима: Проверенный + Автономный + Тест",
        "Журнал сделок с заметками оператора",
        "Walk-forward калибровка с откатом",
        "Telegram + Web-push + Native уведомления",
        "PWA для мобильного через ngrok",
        "Бэктест с CSV-загрузкой истории",
        "Auto-update через GitHub Releases",
    ]
    next_up = [
        "Расширение библиотеки стратегий",
        "Multi-broker одновременно (MT5 + крипта)",
        "Облачная синхронизация Trade Journal",
        "Сценарии входа из ML-моделей",
        "Поддержка macOS (нативный билд)",
        "REST API для внешних клиентов",
    ]

    # Done card
    add_card(s, Inches(0.6), Inches(2.4), Inches(6.0), Inches(4.4),
             border=LAUREL, border_w=1.2)
    add_pill(s, Inches(0.85), Inches(2.6), "ГОТОВО · v1.0.40",
             fill=LAUREL, fg=MARBLE, width=Inches(2.3))
    for i, item in enumerate(done):
        y = Inches(3.15) + i * Inches(0.32)
        check = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(0.95), y + Inches(0.1),
                                    Inches(0.13), Inches(0.13))
        check.fill.solid(); check.fill.fore_color.rgb = LAUREL
        check.line.fill.background()
        add_text(s, Inches(1.2), y, Inches(5.3), Inches(0.3),
                 item, size=10, color=INK)

    # Next card
    add_card(s, Inches(6.85), Inches(2.4), Inches(5.85), Inches(4.4),
             border=GOLD_DEEP, border_w=1.2)
    add_pill(s, Inches(7.1), Inches(2.6), "В РАЗРАБОТКЕ",
             fill=GOLD_DEEP, fg=MARBLE, width=Inches(2.0))
    for i, item in enumerate(next_up):
        y = Inches(3.15) + i * Inches(0.52)
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  Inches(7.2), y + Inches(0.1),
                                  Inches(0.25), Inches(0.18))
        arr.fill.solid(); arr.fill.fore_color.rgb = GOLD_DEEP
        arr.line.fill.background()
        add_text(s, Inches(7.55), y, Inches(5.0), Inches(0.5),
                 item, size=11, color=INK, line_spacing=1.25)

    add_footer(s, page, total)


@slide
def s_thanks(prs, page, total):
    s = prs.slides.add_slide(BLANK); add_bg(s)

    # Decorative side band
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              0, 0, Inches(0.35), SH)
    band.line.fill.background(); band.fill.solid()
    band.fill.fore_color.rgb = GOLD_DEEP
    band.shadow.inherit = False

    add_eyebrow(s, Inches(0.9), Inches(0.9), "BAI Core · Hermes")
    add_text(s, Inches(0.9), Inches(1.5), Inches(11), Inches(1.5),
             "Спасибо.", size=72, bold=True, color=NAVY, font="Cambria")
    add_rule(s, Inches(0.9), Inches(3.0), Inches(2.5), color=GOLD)

    add_text(s, Inches(0.9), Inches(3.25), Inches(11), Inches(1.5),
             "Готовы запустить бот в работу.",
             size=22, color=GOLD_DEEP, font="Cambria")
    add_text(s, Inches(0.9), Inches(4.1), Inches(11), Inches(1.5),
             "Один установщик. Полный контроль. "
             "Объяснимые решения. Производственная стабильность.",
             size=14, color=MUTED, font="Calibri", line_spacing=1.5)

    # Contact card
    add_card(s, Inches(0.9), Inches(5.6), Inches(7.5), Inches(1.2),
             border=GOLD)
    add_text(s, Inches(1.1), Inches(5.75), Inches(7.0), Inches(0.4),
             "baicore.kz", size=14, bold=True, color=NAVY, font="Cambria")
    add_text(s, Inches(1.1), Inches(6.15), Inches(7.0), Inches(0.4),
             "github.com/artsignstudiokz/Hermes · "
             "support@baicore.kz",
             size=11, color=MUTED, font="Consolas")


# ── Render ─────────────────────────────────────────────────────────────


total = len(SLIDES)
for i, builder in enumerate(SLIDES, 1):
    builder(prs, i, total)

prs.save(OUT)
print(f"Saved: {OUT}  ({total} slides)")
